#!/usr/bin/env python3
"""Generate PPT from the Agentic AI Head CPU review report."""

import sys
sys.path.insert(0, ".venv/lib/python3.13/site-packages")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = ACCENT
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, image_path=None, image_left=True):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content area
    if image_path and os.path.exists(image_path):
        if image_left:
            img_left = Inches(0.5)
            text_left = Inches(6.5)
            text_width = Inches(6.3)
        else:
            img_left = Inches(6.5)
            text_left = Inches(0.5)
            text_width = Inches(6.3)
        
        try:
            slide.shapes.add_picture(image_path, img_left, Inches(1.3), width=Inches(5.8))
        except Exception:
            pass
    else:
        text_left = Inches(0.5)
        text_width = Inches(12.3)
    
    content_box = slide.shapes.add_textbox(text_left, Inches(1.3), text_width, Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

def add_two_col_slide(prs, title, left_bullets, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.8))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)
    
    return slide

# ======== SLIDES ========

# Slide 1: Title
add_title_slide(prs,
    "Agentic AI 推理机头 CPU 综述",
    "从 Host 到 Orchestrator\n2026.04.24")

# Slide 2: Background
add_content_slide(prs, "背景：为什么瓶颈正在外溢到 CPU？", [
    "Agentic AI 将推理从『单次请求、连续 decode』转向『多阶段、可暂停、可恢复、可分叉』的复合执行模式",
    "Georgia Tech & Intel (2025-11)：工具处理占端到端延迟的 50%–90.6%",
    "DeepSeek V4 验证趋势：Engram 架构将静态知识检索放在 CPU RAM，动态推理放在 GPU",
    "经济学事实：CPU RAM 每 GB 成本比 GPU HBM 低 10–20 倍",
    "核心判断：机头 CPU 已从传统 host 演化为 inference orchestration layer"
])

# Slide 3: Four Main Lines Overview
add_content_slide(prs, "四条技术主线同时收敛", [
    "1. 算子下发与状态驱动调度：权重量化越激进，『调度墙』越明显",
    "2. KV 卸载与生命周期管理：CPU 内存从 spill 层升级为 warm tier",
    "3. MoE 推理与专家编排：冷专家命中触发同步 CPU→GPU 拷贝",
    "4. PD 分离与跨池编排：CPU 从单节点调度器升级为跨池编排中枢",
    "Agentic workload 的催化下，四条主线形成正反馈"
], image_path="assets/nvidia-dynamo-agentic-kv-readwrite-2026.webp", image_left=False)

# Slide 4: Line 1 - Operator Dispatch
add_content_slide(prs, "主线一：算子下发——从『发命令』到『编排状态机』", [
    "IQ4/FP4 激进量化后，135M 参数模型单次前向发射 301 个 Kernel",
    "纯下发税达 750 μs，占单 Token 时间的 95%",
    "CPU oversubscription 可使 dequeue 延迟放大 19×（vLLM TP=4 场景）",
    "DeepSeek V4 FP4 量化使单 token FLOPs 降至 V3.2 的 27%，下发税占比进一步上升",
    "vLLM V1 通过 Persistent Batch、Numpy 替代 Python Native 重构，吞吐提升 1.7×"
], image_path="assets/extracted/cpu-slowdown-01.png", image_left=False)

# Slide 5: Line 2 - KV Offloading
add_content_slide(prs, "主线二：KV 卸载——从『容量兜底』到『生命周期管理』", [
    "NVIDIA Dynamo：agentic inference cache hit 达 85%–97%，read/write ratio 高达 11.7x",
    "系统价值重心从『写新 KV』转向『保留、路由、预取和恢复旧 KV』",
    "DeepSeek V4：CSA 4x 压缩 + HCA 128x 压缩，KV cache 降至 V3.2 的 10%",
    "NOSA 实现 5.04× 解码吞吐提升；ScoutAttention 通过 Layer-Ahead CPU 预计算减少 bubbles",
    "CXL 内存扩展可将 GPU 需求降低 87%；NVIDIA ICMSP (BlueField-4) 实现 5x token 吞吐"
], image_path="assets/extracted/nosa-01.png", image_left=True)

# Slide 6: Line 3 - MoE
add_content_slide(prs, "主线三：MoE 推理——从『稀疏计算优势』到『Host-side Orchestration 压力』", [
    "DeepSeek V4：1.6T 总参 / 49B 激活参，单节点无法容纳全部专家",
    "冷专家命中触发同步 CPU→GPU 拷贝，CPU 承担路由、权重搬运、拓扑感知三重职责",
    "Speculating Experts：推测预取可将 TPOT 降低 14%",
    "FineMoE：通过 expert map 相似性搜索实现细粒度预取",
    "SpecMoEOff：结合 speculative decoding 实现 2.5× decode 吞吐提升"
], image_path="assets/extracted/spec-experts-01.png", image_left=False)

# Slide 7: Line 4 - PD Disaggregation
add_content_slide(prs, "主线四：PD 分离——从『单节点调度器』到『跨池编排中枢』", [
    "PD 分离已成为几乎每个主要 LLM 服务栈的默认架构",
    "同节点 KV 传输开销 <0.1%，跨节点需 90 Gbps+ 带宽",
    "Agentic batch inference 引入 middle-phase thrashing：异步推进的 agent 其 KV 被 LRU 驱逐",
    "跨池恢复时需反复重算或传输， decode 池长期维持大量并发流 KV 状态",
    "UCCL (NVIDIA) 将跨节点 KV 传输尾延迟降至 7.1%"
], image_path="assets/nvidia-k8s-disagg-serving-2026.webp", image_left=True)

# Slide 8: Real Workloads
add_two_col_slide(prs, "真实工作负载暴露的三项遗漏",
    [
        "高频 prefill 调度",
        "每次工具调用返回触发新的 prefill",
        "每次 subagent 分叉触发 KV 复制",
        "每次聚合触发多路 KV 合并"
    ],
    [
        "多上下文并存管理",
        "OpenClaw 同时维护 10+ 工具调用上下文",
        "Kimi Swarm 多代理并行产生 100+ 活跃 KV 前缀",
        "Claude Code subagents 导致前缀树爆炸"
    ])

# Slide 9: Platform Signals
add_content_slide(prs, "平台信号：硬件路线图围绕 CPU 控制平面收敛", [
    "NVIDIA Vera CPU：88 核 / 1.2 TB/s LPDDR5X / 200 GB/s NVLink-C2C",
    "BlueField-4 STX + ICMSP：5x token 吞吐 / 4x 能效提升，绕过 CPU 瓶颈直接 SSD→GPU KV 卸载",
    "TrendForce：CPU:GPU 配比从 1:8 收敛至 1:1–1:2",
    "Morgan Stanley：DRAM 将取代 HBM 成为 AI 基础设施最紧缺瓶颈，DDR5 价格 Q2 2026 预计涨 50%+",
    "DeepSeek V4 原生华为 Ascend 优化，85%+ 硬件利用率——机头 CPU 优化不限于 x86/Arm"
], image_path="assets/nvidia-vera-cpu-architecture.png", image_left=False)

# Slide 10: Consensus Matrix
add_content_slide(prs, "共识矩阵：四条主线的交互关系", [
    "Kernel Fusion ↗ → 单批次请求数 ↗ → KV 状态数 ↗ → warm tier 压力 ↗",
    "MoE expert 卸载 → CPU RAM 带宽竞争 → 同一节点 KV 预取带宽被挤压",
    "PD 分离 → 跨池 KV 传输需求 ↗ → 机头 CPU 需同时管理本地 warm tier + 远程传输",
    "Agentic 并发度 ↗ → 同时活跃前缀树深度 ↗ → CPU 侧索引结构复杂度指数级上升",
    "优化单一主线可能加剧其他主线瓶颈——需要系统级协同优化"
])

# Slide 11: Conclusion
add_content_slide(prs, "结论：机头 CPU 的选型诊断清单", [
    "核心职责已从 kernel launch 扩展到请求接入、PD 切分、KV 生命周期、专家放置、多代理并发控制",
    "量化越激进，调度墙越明显——『小模型高并发』场景 CPU 选型权重应高于 GPU FLOPs",
    "CPU RAM 容量与带宽成为独立瓶颈：DDR5/LPDDR5X 通道数与容量比核心数更重要",
    "一致性互连（NVLink-C2C / CXL）决定 CPU 能否有效承担 warm tier 角色",
    "机头 CPU 不再是『配一台便宜的就行』，而是 inference factory 的控制平面"
])

# Slide 12: Thank you
add_title_slide(prs,
    "谢谢",
    "完整报告与引用材料详见 GitHub\ngithub.com/hugulas/agentic_cpu_study")

# Save
prs.save('agentic-ai-head-cpu-review.pptx')
print(f"Saved: agentic-ai-head-cpu-review.pptx ({len(prs.slides)} slides)")
